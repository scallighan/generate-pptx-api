from fastapi import FastAPI, Query, Header, Request, Body
from fastapi.responses import FileResponse
from typing import Union, Any
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from jinja2 import Template, exceptions
import json
import httpx
import re
import time

app = FastAPI(
    title="Generate PPTX API",
)

@app.get(
    "/hello",
    tags=["APIs"],
    response_model=dict,
)
async def hello(prompt: Union[str, None] = Query(default="world", max_length=50)):
    print(f"Received prompt: {prompt}")
    return {"message": f"Hello, {prompt}!"}

@app.get(
    "/read",
    tags=["APIs"],
    response_model=dict,
)
async def read_pptx():
    prs = Presentation('/code/app/templates/template.pptx')
    for slide in prs.slides:
        # for shape in slide.shapes:
        #     if shape.has_text_frame:
        #         print(repr(shape.text))
        print(f"Slide id: {slide.slide_id}")
        print(f"Slide Name: {slide.name}")
        print(f"Slide Layout: {slide.slide_layout}")
        print(f"Slide Placeholders:")
        for placeholder in slide.placeholders:
            print(f"    Placeholder name: {placeholder.name}")
            print(f"        type: {placeholder.placeholder_format.type},") 
            print(f"        has_text_frame: {placeholder.has_text_frame},")
            print(f"        text: {placeholder.text}, ")
            print(f"        shape_id: {placeholder.shape_id} ")
            
    return {"message": "PPTX content read successfully."}

async def generate_pptx(data: Any) -> str:
    template = data.get("template", "/code/app/templates/template.pptx")
    prs = Presentation(template)

    for i in range(len(prs.slides)):
        images_index = 0
        tables_index = 0
        slide = prs.slides[i]
        slide_data = data["slides"][i]
        for placeholder in slide.placeholders:
            if placeholder.has_text_frame:
                template = Template(placeholder.text)
                output = template.render(slide_data)
                print(repr(output))
                if output.strip() != "":
                    placeholder.text = output
            if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.TABLE:
                print(f"Found table placeholder: {placeholder.name} | {placeholder.placeholder_format.type}")
                table_arr = slide_data.get("tables", [])
                if len(table_arr) > 0:
                    table_data = table_arr[tables_index]
                    if table_data:
                        rows = len(table_data["rows"]) + 1 # +1 for header
                        cols = len(table_data["headers"])
                        table = placeholder.insert_table(rows, cols).table
                        # Set header
                        for col_index, header in enumerate(table_data["headers"]):
                            table.cell(0, col_index).text = header
                        # Set rows
                        for row_index, row_data in enumerate(table_data["rows"], start=1):
                            for col_index, cell_data in enumerate(row_data):
                                table.cell(row_index, col_index).text = str(cell_data)
                        tables_index += 1

            # Picture placeholder needs to go last as it becomes invalid after insertion    
            elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
                print(f"Found picture placeholder: {placeholder.name} | {placeholder.placeholder_format.type}")
                pictures_arr = slide_data.get("pictures", [])
                if len(pictures_arr) > 0:
                    try:
                        picture_url = pictures_arr[images_index]
                        # regular expression to determine if it is a jpg or png
                        file_ext = re.search(r'\.(jpg|jpeg|png|gif|bmp|tiff)', picture_url, re.IGNORECASE)              
                        img_path = f"/code/app/temp/slide_{i}_image_{images_index}.{file_ext.group(1) if file_ext else 'jpg'}"
                        # Download image
                        print(f"Downloading image from URL: {picture_url}")
                        async with httpx.AsyncClient() as client:
                            async with client.stream("GET", picture_url, follow_redirects=True) as response:
                                response.raise_for_status()  # Raise error for HTTP failures
                                with open(img_path, "wb") as file:
                                    async for chunk in response.aiter_bytes(1024):
                                        file.write(chunk)
                                # Insert picture
                                if(response.status_code == 200):
                                    picture = placeholder.insert_picture(img_path)
                                    images_index += 1
                                    print(f"Inserted picture from {picture_url} [{img_path}] into slide.")
                    except Exception as e:
                        print(f"Error downloading or inserting image: {e}")

            
    output_name = f'generated_presentation_{int(time.time())}.pptx'
    output_path = f'/code/app/output/{output_name}'
    prs.save(output_path)

    #return FileResponse(output_path, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename='generated_presentation.pptx')
    return f"/download?file_path={output_name}" 


@app.get(
    "/download",
    tags=["APIs"],
    response_class=FileResponse,
)
async def download_pptx(file_path: Union[str, None] = Query(default=None, max_length=200)):
    
    
    if file_path is None:
        return {"error": "file_path query parameter is required."}
    modified_file_path = file_path.replace("..", "") if file_path else None
    modified_file_path = f"/code/app/output/{modified_file_path}"
    return FileResponse(modified_file_path, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename=file_path)


# @app.get(
#     "/generate",
#     tags=["APIs"],
# )
async def generate_pptx_static(request: Request, title: Union[str, None] = Query(default="Hello World", max_length=50), content: Union[str, None] = Query(default="Something at nothing", max_length=50)): 
    data = {
        "template": '/code/app/templates/template.pptx',
        "slides": [
            {
                "title": title 
            },
            {
                "agenda": [
                    "Introduction",
                    "Main Topic",
                    "Conclusion"
                ]
            },
            {
                "title": "Introduction - Powerpoint Generation with GenAI",
                "pictures": [
                    "https://learn.microsoft.com/en-us/azure/ai-foundry/how-to/media/managed-virtual-network/diagram-managed-network.png?view=foundry-classic"
                ]
            },
            {
                "sectionname": "Main Topic",
            },
            {

                "title": "We can now generate slides!",
                "content": "We take a template PPTX with placeholders and fill them with data using Jinja2 templating.",
                "pictures": [
                   "https://cdn-dynmedia-1.microsoft.com/is/image/microsoftcorp/Content-Card-Xbox-Controllers-Black-Detail"
                ]
            },
            {
                "title": "Conclusion",
                "content": content,
                "tables": [
                    { 
                        "headers": ["Year", "Earnings", "Profit"],
                        "rows": [
                            ["2023", 30000000, 5000000],
                            ["2024", 25000000, 4000000],
                            ["2025", 35000000, 6000000]
                        ]
                    }
                ]
            },
            {
                "thanks": "Thank you for your attention!",
                "contactinfo": "GenAI PPTX Bot\n- contact@nothing.com\n- 555-123-4567"
            }
        ]
    }
    return await generate_pptx(data)


############################################################

class SlideContent:
    def __init__(self, title: str = None, subtitles: list = None, text: list = None, content: list = None, pictures: list = None, tables: list = None):
        self.title = title
        self.subtitles = subtitles or []
        self.text = text or []
        self.content = content or []
        self.pictures = pictures or []
        self.tables = tables or []
        title_count = 1 if title else 0
        self.weight = title_count + (len(self.subtitles) * 2) + (len(self.text) * 5) + (len(self.content) * 10) + (len(self.pictures) * 100) + (len(self.tables) * 1000)

class SlideLayoutTemplate:
    def __init__(self, layout_name: str, title_count: int = 0, subtitles_count: int = 0, text_count: int = 0, content_count: int = 0, pictures_count: int = 0, tables_count: int = 0):
        self.layout_name = layout_name
        self.title_count = title_count
        self.subtitles_count = subtitles_count
        self.text_count = text_count
        self.content_count = content_count
        self.pictures_count = pictures_count
        self.tables_count = tables_count
        self.weight = title_count + (subtitles_count * 2) + (text_count * 5) + (content_count * 10) + (pictures_count * 100) + (tables_count * 1000)

async def slide_analysis(prs: Presentation) -> dict:
    layout_templates = {}
    for slo in prs.slide_layouts:
        print(f"{[slo.name]} Layout Analysis:")
        title_count = 0
        subtitles_count = 0
        text_count = 0
        content_count = 0
        pictures_count = 0
        tables_count = 0
        for placeholder in slo.placeholders:
            print(f"    Placeholder name: {placeholder.name}")
            print(f"        type: {placeholder.placeholder_format.type},")
            if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE or placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.CENTER_TITLE:
                title_count += 1
            elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.BODY:
                text_count += 1
            elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE:
                subtitles_count += 1
            elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.OBJECT:
                content_count += 1
            elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
                pictures_count += 1
            elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.TABLE:
                tables_count += 1
        layout_template = SlideLayoutTemplate(
            layout_name=slo.name,
            title_count=title_count,
            subtitles_count=subtitles_count,
            text_count=text_count,
            content_count=content_count,
            pictures_count=pictures_count,
            tables_count=tables_count
        )
        layout_templates[slo.name.lower()] = layout_template
    return layout_templates       

async def _populate_slide(slide, slide_content: SlideContent):
    subtitles_count =0
    text_count = 0
    content_count = 0
    pictures_count = 0
    tables_count = 0

    for placeholder in slide.placeholders:
        if (placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE or placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.CENTER_TITLE) and slide_content.title:
            placeholder.text = slide_content.title
        elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE and len(slide_content.subtitles) > 0:
            placeholder.text = slide_content.subtitles[subtitles_count]
            subtitles_count += 1
        elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.BODY and len(slide_content.text) > 0:
            placeholder.text = slide_content.text[text_count]
            text_count += 1
        elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.OBJECT and len(slide_content.content) > 0:
            placeholder.text = slide_content.content[content_count]
            content_count += 1
        elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE and len(slide_content.pictures) > 0:
            try:
                picture_url = slide_content.pictures[pictures_count]
                pictures_count += 1
                file_ext = re.search(r'\.(jpg|jpeg|png|gif|bmp|tiff)', picture_url, re.IGNORECASE)              
                img_path = f"/code/app/temp/temp_image.{file_ext.group(1) if file_ext else 'jpg'}"
                # Download image
                print(f"Downloading image from URL: {picture_url}")
                async with httpx.AsyncClient() as client:
                    async with client.stream("GET", picture_url, follow_redirects=True) as response:
                        response.raise_for_status()  # Raise error for HTTP failures
                        with open(img_path, "wb") as file:
                            async for chunk in response.aiter_bytes(1024):
                                file.write(chunk)
                        # Insert picture
                        if(response.status_code == 200):
                            picture = placeholder.insert_picture(img_path)
                            print(f"Inserted picture from {picture_url} [{img_path}] into slide.")
            except Exception as e:
                print(f"Error downloading or inserting image: {e}")
        elif placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.TABLE and len(slide_content.tables) > 0:
            try:
                table_data = slide_content.tables[tables_count]
                tables_count += 1
                if table_data:
                    rows = len(table_data["rows"]) + 1 # +1 for header
                    cols = len(table_data["headers"])
                    table = placeholder.insert_table(rows, cols).table
                    # Set header
                    for col_index, header in enumerate(table_data["headers"]):
                        table.cell(0, col_index).text = header
                    # Set rows
                    for row_index, row_data in enumerate(table_data["rows"], start=1):
                        for col_index, cell_data in enumerate(row_data):
                            table.cell(row_index, col_index).text = str(cell_data)
                    print(f"Inserted table into slide.")
            except Exception as e:
                print(f"Error inserting table: {e}")

async def generate_pptx_v2(data: Any) -> str:
    template = data.get("template", "/code/app/templates/template.pptx")
    prs = Presentation(template)

    slide_layout_analysis = await slide_analysis(prs)
    print(f"Slide Layout Analysis Results: {slide_layout_analysis}")
    slide_data = data.get("slides", [])
    
    for s in slide_data:
        print(f"Processing slide data: {s}")
        slide_content = SlideContent(
            title=s.get("title"),
            subtitles=s.get("subtitles", []),
            text=s.get("text", []),
            content=s.get("content", []),
            pictures=s.get("pictures", []),
            tables=s.get("tables", [])
        )
        best_match = None
        if slide_layout_analysis.get(slide_content.title.lower()):
            print(f"Matched layout: {slide_layout_analysis[slide_content.title.lower()].layout_name} for slide content with weight {slide_content.weight}")
            best_match = slide_layout_analysis[slide_content.title.lower()]
        else:
            # find the best matching weight by comparing slide_content.weight to layout_template.weight
            
            smallest_diff = float('inf')
            for layout_name, layout_template in slide_layout_analysis.items():
                weight_diff = abs(slide_content.weight - layout_template.weight)
                if weight_diff < smallest_diff:
                    smallest_diff = weight_diff
                    best_match = layout_template
        if best_match:
            print(f"Best matched layout: {best_match.layout_name} ({best_match.weight}) for slide content with weight {slide_content.weight}")
            current_slide = prs.slides.add_slide(prs.slide_layouts.get_by_name(best_match.layout_name))
            await _populate_slide(current_slide, slide_content)
        else:
            print(f"No matching layout found for slide content with weight {slide_content.weight}")

        
    output_name = f'generated_presentation_{int(time.time())}.pptx'
    output_path = f'/code/app/output/{output_name}'
    prs.save(output_path)
    return f"/download?file_path={output_name}" 

# @app.post(
#     "/dynamic",
#     tags=["APIs"],
#     description="New functionality testing"
# )
# async def dynamic_function(request: Request, body: dict = Body(...)):
#     downloadpath = await generate_pptx_v2(body)
#     return { "download_url": f"{'http' if request.url.port else 'https'}://{request.url.hostname}:{request.url.port if request.url.port else '443'}{downloadpath}" }


@app.post(
    "/generate",
    tags=["APIs"],
    description="Generate a PPTX presentation based on the provided JSON structure."
)
async def generate_pptx_dynamic(request: Request, body: dict = Body(...)): 
    downloadpath = await generate_pptx_v2(body)
    return { "download_url": f"{'http' if request.url.port else 'https'}://{request.url.hostname}:{request.url.port if request.url.port else '443'}{downloadpath}" }

